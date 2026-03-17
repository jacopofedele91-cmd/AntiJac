import os
from docx import Document
from pptx import Presentation

def extract_docx(file_path):
    print(f"Extracting {file_path}")
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip() != ""])

def extract_pptx(file_path):
    print(f"Extracting {file_path}")
    prs = Presentation(file_path)
    text_runs: list[str] = []
    for slide_i, slide in enumerate(prs.slides):
        text_runs.append(f"--- SLIDE {slide_i+1} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        text_runs.append(run.text.strip())
    return "\n".join(text_runs)

def main():
    desktop = r"c:\Users\Acer\OneDrive\Desktop\AntiJac"
    tmp_dir = os.path.join(desktop, ".tmp")
    os.makedirs(tmp_dir, exist_ok=True)

    output = []

    files_to_read = [
        ("QUESTIONARIO INIZIALE.docx", extract_docx),
        ("TABELLA REGOLE_PRODOTTO PER TE.docx", extract_docx),
        ("FRONT END APP.pptx", extract_pptx)
    ]

    for fname, extractor in files_to_read:
        fpath = os.path.join(desktop, fname)
        if os.path.exists(fpath):
            try:
                text = extractor(fpath)
                output.append(f"=== {fname} ===\n{text}\n")
            except Exception as e:
                output.append(f"=== {fname} ===\nERROR: {e}\n")
        else:
            output.append(f"=== {fname} ===\nNOT FOUND\n")

    with open(os.path.join(tmp_dir, "extracted_docs.txt"), "w", encoding="utf-8") as f:
        f.write("\n".join(output))
    print("Extraction complete.")

if __name__ == "__main__":
    main()
